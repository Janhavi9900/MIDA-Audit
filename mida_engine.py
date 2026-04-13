"""
mida_engine.py
MIDA audit engine for Benchling Canvas App.

Schema: assaysch_Q0VwMVfrjo (MIDA Result table 1)
Fields:
  samples  (entry_link) — linked sample entity
  status   (text)       — MIDA writes VERIFIED / JUSTIFIED / FAIL here
  docs     (blob_link)  — linked SOP / rule document per sample row

Flow per row:
  1. Get sample name from samples field
  2. Get SOP text from docs blob
  3. Run Gemini AI clash detection (same logic as Streamlit app.py)
  4. Run Gemini final audit with AND logic
  5. Write VERIFIED / JUSTIFIED / FAIL to status field
"""

import os
import io
import json
import base64
import requests
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel("models/gemini-2.5-flash")

BENCHLING_URL  = os.getenv("BENCHLING_TENANT_URL")
BENCHLING_KEY  = os.getenv("BENCHLING_API_KEY")
RESULTS_SCHEMA = "assaysch_Q0VwMVfrjo"


def _api_headers():
    creds = base64.b64encode(f"{BENCHLING_KEY}:".encode()).decode()
    return {
        "Authorization": f"Basic {creds}",
        "Content-Type":  "application/json"
    }


def fetch_blob_content(blob_id: str) -> bytes:
    """Download raw bytes of a blob."""
    resp = requests.get(
        f"{BENCHLING_URL}/api/v2/blobs/{blob_id}/download",
        headers=_api_headers(),
        timeout=60
    )
    if resp.status_code == 200:
        return resp.content
    print(f"  Could not fetch blob {blob_id}: {resp.status_code}")
    return b""


def blob_to_text(content: bytes, filename: str) -> str:
    """Convert blob bytes to text based on file type."""
    fname = filename.lower()

    try:
        if fname.endswith(".pdf"):
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            return "\n".join(p.extract_text() or "" for p in reader.pages)

        elif fname.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)

        elif fname.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            return "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

        elif fname.endswith((".xlsx", ".xls")):
            import pandas as pd
            df = pd.read_excel(io.BytesIO(content))
            return df.to_csv(index=False)

        elif fname.endswith(".csv"):
            return content.decode("utf-8", errors="ignore")

        else:
            # txt, md, or unknown — try plain text
            return content.decode("utf-8", errors="ignore")

    except Exception as e:
        print(f"  Could not parse {filename}: {e}")
        return content.decode("utf-8", errors="ignore")


def get_sample_assay_data(sample_id: str, entry_id: str) -> str:
    """
    Fetch assay result data for a specific sample from the entry.
    Returns CSV string of matching assay results for Gemini to audit.
    """
    resp = requests.get(
        f"{BENCHLING_URL}/api/v2/assay-results",
        headers=_api_headers(),
        params={"entryIds": entry_id},
        timeout=30
    )
    if resp.status_code != 200:
        return ""

    rows = resp.json().get("assayResults", [])
    matched = []
    for row in rows:
        fields = row.get("fields", {})
        # Look for any field containing the sample ID
        for field_name, field_data in fields.items():
            if isinstance(field_data, dict):
                val = field_data.get("value", "")
                if isinstance(val, dict):
                    val = val.get("name", "") or val.get("id", "")
                val = str(val)
            else:
                val = str(field_data)

            if sample_id.lower() in val.lower() or val.lower() in sample_id.lower():
                # Flatten this row's fields to a dict
                flat = {}
                for fn, fd in fields.items():
                    if isinstance(fd, dict):
                        v = fd.get("value", "")
                        if isinstance(v, dict):
                            v = v.get("name", "") or v.get("id", "")
                    else:
                        v = fd
                    flat[fn] = v
                matched.append(flat)
                break

    if not matched:
        return f"Sample: {sample_id}\nNo additional assay result data found."

    import pandas as pd
    return pd.DataFrame(matched).to_csv(index=False)


def audit_single_sample(sample_name: str, sample_data: str, rule_text: str) -> dict:
    """
    Run Gemini audit for a single sample against its rule document.
    Mirrors the 2-step logic from app.py (clash detection + final audit).
    Returns {status, comment}.
    """

    # Step 1 — Clash detection (same as app.py clash_prompt)
    clash_prompt = f"""
You are a senior lab compliance auditor. Analyse lab results against rule documents.

SAMPLE: {sample_name}

LAB DATA:
{sample_data}

RULE DOCUMENT:
{rule_text}

INSTRUCTIONS:
1. Check each value in the lab data against the rule document.
2. A CLASH exists ONLY when a value PASSES in one interpretation but FAILS in another.
3. Extract EXACT thresholds with numbers and units from the document.

Return ONLY a JSON array (empty [] if no clashes):
[{{
  "parameter": "exact parameter name",
  "cell_value": "actual value",
  "clash_summary": "one sentence",
  "groups": [
    {{"stance": "PASS", "rule": "...", "threshold": "..."}},
    {{"stance": "FAIL", "rule": "...", "threshold": "..."}}
  ]
}}]
No markdown — ONLY the JSON array.
"""

    try:
        clash_resp = model.generate_content(clash_prompt)
        clash_raw  = clash_resp.text.replace("```json","").replace("```","").strip()
        clashes    = json.loads(clash_raw)
    except Exception as e:
        print(f"  Clash detection error: {e}")
        clashes = []

    # Step 2 — Final audit with AND logic (same as app.py final_prompt)
    final_prompt = f"""
You are a Lead Lab Auditor generating a final compliance decision.

SAMPLE: {sample_name}

LAB DATA:
{sample_data}

RULE DOCUMENT:
{rule_text}

DETECTED CLASHES:
{json.dumps(clashes)}

AUDIT RULES:
1. AND LOGIC — sample PASSES only if ALL parameters pass. One fail = whole sample FAILS.
2. STATUS VALUES:
   - VERIFIED  — all parameters pass the rule document
   - JUSTIFIED — parameter failed rule doc but a clash shows it passes another interpretation
   - FAIL      — one or more parameters fail with no justification
3. Comment must list every parameter checked with its value and threshold.

Return ONLY a JSON object — no markdown:
{{"status": "VERIFIED|JUSTIFIED|FAIL", "comment": "explanation max 150 chars"}}
"""

    try:
        final_resp = model.generate_content(final_prompt)
        final_raw  = final_resp.text.replace("```json","").replace("```","").strip()
        result     = json.loads(final_raw)
        return {
            "status":  result.get("status", "FAIL"),
            "comment": result.get("comment", "")
        }
    except Exception as e:
        print(f"  Final audit error: {e}")
        return {"status": "FAIL", "comment": f"Audit error: {str(e)[:80]}"}


def run_audit_on_entry(entry_id: str) -> list:
    """
    Read MIDA Result table rows for this entry.
    For each row: get sample + doc → run Gemini audit.
    Returns list of {result_id, sample_id, status, comment}.
    """
    if not entry_id:
        raise Exception("No entry ID provided.")

    headers = _api_headers()

    # Fetch MIDA result table rows for this entry
    print(f"  Fetching MIDA result rows for entry: {entry_id}")
    resp = requests.get(
        f"{BENCHLING_URL}/api/v2/assay-results",
        headers=headers,
        params={
            "schemaId": RESULTS_SCHEMA,
            "entryIds": entry_id,
        },
        timeout=30
    )

    if resp.status_code != 200:
        raise Exception(
            f"Could not fetch MIDA result rows: {resp.status_code} {resp.text[:300]}"
        )

    rows = resp.json().get("assayResults", [])
    print(f"  Found {len(rows)} MIDA result rows")

    if not rows:
        raise Exception(
            "No rows found in the MIDA Result table for this entry. "
            "Please add rows with Samples and Docs filled in."
        )

    audit_results = []

    for row in rows:
        fields    = row.get("fields", {})
        result_id = row.get("id")

        # ── Get sample name ──────────────────────────────────────────
        samples_field = fields.get("samples", {})
        if isinstance(samples_field, dict):
            val = samples_field.get("value")
            if isinstance(val, dict):
                sample_name = val.get("name") or val.get("id") or str(val)
            elif val:
                sample_name = str(val)
            else:
                sample_name = f"Row-{result_id}"
        else:
            sample_name = f"Row-{result_id}"

        print(f"  Processing sample: {sample_name}")

        # ── Get rule document from docs field ────────────────────────
        docs_field = fields.get("docs", {})
        rule_text  = ""
        doc_name   = "Rule Document"

        if isinstance(docs_field, dict):
            blob_val = docs_field.get("value")
            if isinstance(blob_val, dict):
                blob_id  = blob_val.get("id") or blob_val.get("blobId")
                doc_name = blob_val.get("name", "document")
            elif blob_val and isinstance(blob_val, str):
                blob_id  = blob_val
            else:
                blob_id = None

            if blob_id:
                print(f"  Fetching doc blob: {blob_id}")
                content  = fetch_blob_content(blob_id)
                rule_text = blob_to_text(content, doc_name)
                print(f"  Rule doc fetched: {len(rule_text)} chars")

        if not rule_text:
            print(f"  No rule doc found for {sample_name} — will audit without rules")
            rule_text = "No rule document provided. Flag as JUSTIFIED."

        # ── Get sample assay data from other result tables ───────────
        sample_data = get_sample_assay_data(sample_name, entry_id)
        if not sample_data:
            sample_data = f"Sample: {sample_name}\nNo assay data found in entry."

        # ── Run Gemini audit ─────────────────────────────────────────
        print(f"  Running Gemini audit for: {sample_name}")
        result = audit_single_sample(sample_name, sample_data, rule_text)
        print(f"  Result: {result['status']} — {result['comment'][:60]}")

        audit_results.append({
            "result_id": result_id,
            "sample_id": sample_name,
            "status":    result["status"],
            "comment":   result["comment"],
        })

    return audit_results


def write_results_to_benchling(entry_id: str, audit_results: list) -> int:
    """
    Write audit status to the 'status' field in each MIDA Result table row.
    Returns count of rows successfully updated.
    """
    if not audit_results:
        return 0

    headers       = _api_headers()
    updated_count = 0

    for result in audit_results:
        result_id = result.get("result_id")
        status    = result.get("status", "")
        comment   = result.get("comment", "")

        if not result_id:
            print(f"  No result_id for {result.get('sample_id')} — skipping")
            continue

        # Write status + comment to status field
        status_value = f"{status} — {comment[:100]}" if comment else status

        print(f"  Writing {result_id} → {status_value[:60]}")
        update_resp = requests.patch(
            f"{BENCHLING_URL}/api/v2/assay-results/{result_id}",
            headers=headers,
            json={
                "fields": {
                    "status": {"value": status_value}
                }
            },
            timeout=15
        )
        print(f"  Response: {update_resp.status_code}")
        if update_resp.status_code in (200, 201):
            updated_count += 1
        else:
            print(f"  Error: {update_resp.text[:200]}")

    return updated_count